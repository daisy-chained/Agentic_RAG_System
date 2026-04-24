import grpc
import asyncio
import time
import os
from dotenv import load_dotenv

from opentelemetry.exporter.otlp.proto.grpc.trace_exporter import OTLPSpanExporter
from opentelemetry.sdk.trace import TracerProvider
from opentelemetry.sdk.trace.export import SimpleSpanProcessor
from opentelemetry import trace
from openinference.instrumentation.langchain import LangChainInstrumentor

from langchain_ollama import ChatOllama
from langchain_core.messages import SystemMessage, HumanMessage, AIMessage

import ai_service_pb2_grpc as pb2_grpc
import ai_service_pb2 as pb2
import tempfile
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
    CSVLoader,
    BSHTMLLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter

import rag
import reranker

# ── Observability Tracing (Arize Phoenix) ─────────────────────────────────────
trace.set_tracer_provider(TracerProvider())
trace.get_tracer_provider().add_span_processor(
    SimpleSpanProcessor(OTLPSpanExporter(endpoint="http://localhost:4317"))
)
LangChainInstrumentor().instrument()

# ── Ollama Configuration ──────────────────────────────────────────────────────
# Reads from .env locally; Docker Compose injects these as environment variables.
load_dotenv()

_OLLAMA_HOST    = os.getenv("OLLAMA_HOST",    "http://localhost:11434")
_OLLAMA_MODEL   = os.getenv("OLLAMA_MODEL",   "llama3.2")
_OLLAMA_NUM_CTX = int(os.getenv("OLLAMA_NUM_CTX", "4096"))

_SYSTEM_PROMPT = (
    "You are an expert technical AI assistant embedded in a "
    "Retrieval-Augmented Generation (RAG) system. Answer questions "
    "accurately and concisely. When context is provided, ground your "
    "answers in that context. Acknowledge clearly when you are uncertain."
)

_llm = ChatOllama(model=_OLLAMA_MODEL, base_url=_OLLAMA_HOST, num_ctx=_OLLAMA_NUM_CTX)

print(f"✓ Inference Engine configured — ollama host={_OLLAMA_HOST} model={_OLLAMA_MODEL} num_ctx={_OLLAMA_NUM_CTX}")


# ── Custom document loaders (pure-Python, no system-level dependencies) ──────

class _XlsxLoader:
    """Load .xlsx spreadsheets using openpyxl (pure Python)."""

    def __init__(self, file_path: str):
        self._path = file_path

    def load(self):
        import openpyxl
        from langchain_core.documents import Document

        wb = openpyxl.load_workbook(self._path, read_only=True, data_only=True)
        docs = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_str = "\t".join("" if v is None else str(v) for v in row)
                if row_str.strip():
                    rows.append(row_str)
            if rows:
                docs.append(Document(
                    page_content="\n".join(rows),
                    metadata={"sheet": sheet_name},
                ))
        wb.close()
        return docs


class _XlsLoader:
    """Load legacy .xls spreadsheets using xlrd (pure Python)."""

    def __init__(self, file_path: str):
        self._path = file_path

    def load(self):
        import xlrd
        from langchain_core.documents import Document

        wb = xlrd.open_workbook(self._path)
        docs = []
        for sheet in wb.sheets():
            rows = []
            for rx in range(sheet.nrows):
                row_str = "\t".join(
                    str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)
                )
                if row_str.strip():
                    rows.append(row_str)
            if rows:
                docs.append(Document(
                    page_content="\n".join(rows),
                    metadata={"sheet": sheet.name},
                ))
        return docs


class _PptxLoader:
    """Load .pptx presentations using python-pptx (pure Python)."""

    def __init__(self, file_path: str):
        self._path = file_path

    def load(self):
        from pptx import Presentation
        from langchain_core.documents import Document

        prs = Presentation(self._path)
        docs = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame
            ]
            content = "\n".join(texts).strip()
            if content:
                docs.append(Document(
                    page_content=content,
                    metadata={"slide": i},
                ))
        return docs


class _EpubLoader:
    """Load .epub e-books using ebooklib + BeautifulSoup (pure Python)."""

    def __init__(self, file_path: str):
        self._path = file_path

    def load(self):
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        from langchain_core.documents import Document

        book = epub.read_epub(self._path)
        docs = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text:
                docs.append(Document(
                    page_content=text,
                    metadata={"item": item.get_name()},
                ))
        return docs


# ── gRPC Servicer ─────────────────────────────────────────────────────────────
class AiAgentServicer(pb2_grpc.AiAgentServiceServicer):

    async def ProcessQuery(self, request, context):
        start_time = time.monotonic()
        print(f"[ProcessQuery] userId={request.user_id} query='{request.query[:80]}...'")

        # 1. RAG Retrieval via Qdrant → cross-encoder reranking
        retrieved_docs = []
        try:
            # Fetch a wider candidate set so the reranker has more to work with
            candidates = await asyncio.to_thread(
                rag.vector_store.similarity_search_with_score,
                query=request.query,
                k=10
            )
            # Rerank candidates with the cross-encoder and keep the top results
            top_results = await asyncio.to_thread(
                reranker.rerank,
                request.query,
                candidates,
            )
            context_texts = []
            for doc, score in top_results:
                source = doc.metadata.get("source", "Unknown")
                retrieved_docs.append(source)
                context_texts.append(f"--- Excerpt from {source} ---\n{doc.page_content}\n")
            
            rag_context = "\n".join(context_texts)
        except Exception as e:
            print(f"WARN: RAG retrieval failed: {e}")
            rag_context = ""

        # Build the full message list for Ollama
        system_content = _SYSTEM_PROMPT
        if rag_context:
            system_content += f"\n\nGROUND TRUTH CONTEXT to base your answer on:\n{rag_context}"

        messages = [SystemMessage(content=system_content)]

        for item in request.session_history:
            if item.startswith("user: "):
                messages.append(HumanMessage(content=item[6:]))
            elif item.startswith("assistant: "):
                messages.append(AIMessage(content=item[11:]))

        messages.append(HumanMessage(content=request.query))

        try:
            # Generate the draft answer asynchronously
            response = await _llm.ainvoke(messages)
            answer = response.content

            print(f"[Reflection] Triggering evaluation for query '{request.query[:30]}...'")
            reflection_prompt = [
                SystemMessage(content="You are a strict grading evaluator determining whether a Draft Answer is factually grounded in the provided Context. "
                                      "RULES:\n"
                                      "1. If the query cannot be answered using the Context, the Draft Answer MUST clearly state that it does not know based on the context. If it hallucinated an answer anyway, output a REWRITTEN response stating that the information is missing from the context.\n"
                                      "2. If the Draft Answer contains factual claims that contradict or are missing from the Context, output a REWRITTEN response correcting it based ONLY on the Context.\n"
                                      "3. If the Draft Answer is perfectly correct and grounded, you must output exactly and only the word 'APPROVE'.\n"
                                      "4. When outputting a REWRITTEN response, DO NOT include meta-commentary, DO NOT explain why you rewrote it, and DO NOT use prefixes like 'REWRITTEN correct answer:'. Output ONLY the exact text the user should see as the final answer."),
                HumanMessage(content=f"QUERY: {request.query}\n\nCONTEXT: {rag_context}\n\nDRAFT ANSWER: {answer}")
            ]
            
            reflection_response = await _llm.ainvoke(reflection_prompt)
            reflection_text = reflection_response.content.strip()

            if reflection_text.upper() != "APPROVE" and len(reflection_text) > 5:
                print(f"[Reflection] Hallucination detected! Overriding draft answer.")
                answer = f"[Reflexion Self-Corrected]\n{reflection_text}"
            else:
                print(f"[Reflection] Approved.")

        except Exception as e:
            print(f"ERROR: Ollama call failed: {e}")
            context.set_code(grpc.StatusCode.INTERNAL)
            context.set_details(f"Local LLM error: {str(e)}")
            return pb2.AgentResponse()

        latency_ms = int((time.monotonic() - start_time) * 1000)
        print(f"[ProcessQuery] completed in {latency_ms}ms")

        return pb2.AgentResponse(
            answer=answer,
            source_documents=list(set(retrieved_docs)), # unique sources
            confidence_score=0.9, # To dynamically calculate based on distance in Phase 4
            latency_ms=latency_ms,
        )

    async def IndexDocument(self, request, context):
        start_time = time.monotonic()
        print(f"[IndexDocument] Received {request.filename} ({len(request.file_content)} bytes)")

        # 1. Write bytes to a temporary file
        fd, temp_path = tempfile.mkstemp(suffix=f"_{request.filename}")
        try:
            with os.fdopen(fd, 'wb') as f:
                f.write(request.file_content)
            
            # 2. Parse Document — pick loader by file extension
            ext = os.path.splitext(request.filename.lower())[1]
            _loader_map = {
                ".pdf":  PyPDFLoader,
                ".docx": Docx2txtLoader,
                ".xlsx": _XlsxLoader,
                ".xls":  _XlsLoader,
                ".pptx": _PptxLoader,
                ".csv":  CSVLoader,
                ".md":   TextLoader,
                ".html": BSHTMLLoader,
                ".htm":  BSHTMLLoader,
                ".epub": _EpubLoader,
                ".txt":  TextLoader,
            }
            loader = _loader_map.get(ext, TextLoader)(temp_path)
            
            docs = loader.load()

            # Assign proper source metadata
            for doc in docs:
                doc.metadata["source"] = request.filename
                doc.metadata["document_id"] = request.document_id
                doc.metadata["uploaded_by"] = request.user_id

            # 3. Chunking
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200
            )
            chunks = text_splitter.split_documents(docs)

            # 4. Upsert vector embeddings synchronously in thread
            await asyncio.to_thread(
                rag.vector_store.add_documents,
                documents=chunks
            )

            chunk_count = len(chunks)
            print(f"[IndexDocument] Successfully indexed {chunk_count} chunks for {request.filename}")

            return pb2.IndexResponse(
                chunk_count=chunk_count,
                status="SUCCESS"
            )

        except Exception as e:
            print(f"ERROR: Indexing failed -> {e}")
            context.set_code(grpc.StatusCode.INTERNAL)
            context.set_details(str(e))
            return pb2.IndexResponse(status="FAILED")
        finally:
            os.remove(temp_path)


# ── Server ────────────────────────────────────────────────────────────────────
async def serve():
    server = grpc.aio.server()
    pb2_grpc.add_AiAgentServiceServicer_to_server(AiAgentServicer(), server)
    listen_addr = "[::]:50051"
    server.add_insecure_port(listen_addr)
    print(f"✓ Inference Engine gRPC server listening on {listen_addr}")
    await server.start()
    await server.wait_for_termination()


if __name__ == "__main__":
    asyncio.run(serve())